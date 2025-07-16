#!/usr/bin/env python3
"""
简易Office MCP服务器
支持Word、Excel、PPT文档的基本操作
"""

import asyncio
import json
import os
from pathlib import Path
from typing import Any, Dict, List, Optional

# MCP相关导入
from mcp.server import Server
from mcp.server.models import InitializationOptions
from mcp.server.stdio import stdio_server
from mcp.types import (
    CallToolRequest,
    CallToolResult,
    ListToolsRequest,
    ListToolsResult,
    Tool,
    TextContent,
    ImageContent,
    EmbeddedResource
)

# Office文档处理库
try:
    from docx import Document
    from docx.shared import Inches
    import openpyxl
    from openpyxl.styles import Font, PatternFill
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    print(f"请安装所需依赖: pip install python-docx openpyxl python-pptx")
    raise e

class OfficeManager:
    """Office文档管理器"""
    
    def __init__(self, workspace_dir: str = "./office_workspace"):
        self.workspace_dir = Path(workspace_dir)
        self.workspace_dir.mkdir(exist_ok=True)
    
    def get_file_path(self, filename: str) -> Path:
        """获取文件完整路径"""
        return self.workspace_dir / filename
    
    # Word文档操作
    def create_word_document(self, filename: str, title: str = "", content: str = "") -> str:
        """创建Word文档"""
        try:
            doc = Document()
            
            if title:
                doc.add_heading(title, 0)
            
            if content:
                doc.add_paragraph(content)
            
            file_path = self.get_file_path(filename)
            doc.save(file_path)
            
            return f"Word文档 '{filename}' 创建成功，保存在 {file_path}"
        except Exception as e:
            return f"创建Word文档失败: {str(e)}"
    
    def add_word_content(self, filename: str, content: str, heading_level: int = 0) -> str:
        """向Word文档添加内容"""
        try:
            file_path = self.get_file_path(filename)
            
            if not file_path.exists():
                return f"文件 '{filename}' 不存在"
            
            doc = Document(file_path)
            
            if heading_level > 0:
                doc.add_heading(content, heading_level)
            else:
                doc.add_paragraph(content)
            
            doc.save(file_path)
            
            return f"内容已添加到 '{filename}'"
        except Exception as e:
            return f"添加内容失败: {str(e)}"
    
    def read_word_document(self, filename: str) -> str:
        """读取Word文档内容"""
        try:
            file_path = self.get_file_path(filename)
            
            if not file_path.exists():
                return f"文件 '{filename}' 不存在"
            
            doc = Document(file_path)
            content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            
            return "\n".join(content)
        except Exception as e:
            return f"读取Word文档失败: {str(e)}"
    
    # Excel操作
    def create_excel_workbook(self, filename: str, sheet_name: str = "Sheet1") -> str:
        """创建Excel工作簿"""
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet_name
            
            file_path = self.get_file_path(filename)
            wb.save(file_path)
            
            return f"Excel工作簿 '{filename}' 创建成功，保存在 {file_path}"
        except Exception as e:
            return f"创建Excel工作簿失败: {str(e)}"
    
    def write_excel_data(self, filename: str, sheet_name: str, data: List[List[Any]], 
                        start_row: int = 1, start_col: int = 1) -> str:
        """向Excel写入数据"""
        try:
            file_path = self.get_file_path(filename)
            
            if not file_path.exists():
                return f"文件 '{filename}' 不存在"
            
            wb = openpyxl.load_workbook(file_path)
            
            if sheet_name not in wb.sheetnames:
                wb.create_sheet(sheet_name)
            
            ws = wb[sheet_name]
            
            for row_idx, row_data in enumerate(data):
                for col_idx, cell_value in enumerate(row_data):
                    ws.cell(row=start_row + row_idx, 
                           column=start_col + col_idx, 
                           value=cell_value)
            
            wb.save(file_path)
            
            return f"数据已写入 '{filename}' 的 '{sheet_name}' 工作表"
        except Exception as e:
            return f"写入Excel数据失败: {str(e)}"
    
    def read_excel_data(self, filename: str, sheet_name: str = None) -> str:
        """读取Excel数据"""
        try:
            file_path = self.get_file_path(filename)
            
            if not file_path.exists():
                return f"文件 '{filename}' 不存在"
            
            wb = openpyxl.load_workbook(file_path)
            
            if sheet_name is None:
                sheet_name = wb.active.title
            
            if sheet_name not in wb.sheetnames:
                return f"工作表 '{sheet_name}' 不存在"
            
            ws = wb[sheet_name]
            data = []
            
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    data.append(list(row))
            
            return json.dumps(data, ensure_ascii=False, indent=2)
        except Exception as e:
            return f"读取Excel数据失败: {str(e)}"
    
    # PowerPoint操作
    def create_ppt_presentation(self, filename: str, title: str = "") -> str:
        """创建PPT演示文稿"""
        try:
            prs = Presentation()
            
            if title:
                # 添加标题幻灯片
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                slide.shapes.title.text = title
            
            file_path = self.get_file_path(filename)
            prs.save(file_path)
            
            return f"PPT演示文稿 '{filename}' 创建成功，保存在 {file_path}"
        except Exception as e:
            return f"创建PPT演示文稿失败: {str(e)}"
    
    def add_ppt_slide(self, filename: str, title: str, content: str = "", 
                     layout_index: int = 1) -> str:
        """向PPT添加幻灯片"""
        try:
            file_path = self.get_file_path(filename)
            
            if not file_path.exists():
                return f"文件 '{filename}' 不存在"
            
            prs = Presentation(file_path)
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            
            # 设置标题
            if slide.shapes.title:
                slide.shapes.title.text = title
            
            # 设置内容
            if content and len(slide.placeholders) > 1:
                slide.placeholders[1].text = content
            
            prs.save(file_path)
            
            return f"幻灯片已添加到 '{filename}'"
        except Exception as e:
            return f"添加幻灯片失败: {str(e)}"
    
    def get_ppt_info(self, filename: str) -> str:
        """获取PPT信息"""
        try:
            file_path = self.get_file_path(filename)
            
            if not file_path.exists():
                return f"文件 '{filename}' 不存在"
            
            prs = Presentation(file_path)
            
            info = {
                "文件名": filename,
                "幻灯片数量": len(prs.slides),
                "幻灯片标题": []
            }
            
            for i, slide in enumerate(prs.slides):
                title = "无标题"
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text
                info["幻灯片标题"].append(f"第{i+1}页: {title}")
            
            return json.dumps(info, ensure_ascii=False, indent=2)
        except Exception as e:
            return f"获取PPT信息失败: {str(e)}"
    
    def list_files(self) -> str:
        """列出工作区所有文件"""
        try:
            files = []
            for file_path in self.workspace_dir.iterdir():
                if file_path.is_file():
                    files.append({
                        "文件名": file_path.name,
                        "大小": f"{file_path.stat().st_size} bytes",
                        "类型": file_path.suffix
                    })
            
            return json.dumps(files, ensure_ascii=False, indent=2)
        except Exception as e:
            return f"列出文件失败: {str(e)}"

# 创建MCP服务器
server = Server("office-manager")
office_manager = OfficeManager()

@server.list_tools()
async def handle_list_tools() -> List[Tool]:
    """列出所有可用工具"""
    return [
        # Word工具
        Tool(
            name="create_word_document",
            description="创建新的Word文档",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "文件名（包含.docx扩展名）"},
                    "title": {"type": "string", "description": "文档标题（可选）"},
                    "content": {"type": "string", "description": "初始内容（可选）"}
                },
                "required": ["filename"]
            }
        ),
        Tool(
            name="add_word_content",
            description="向Word文档添加内容",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "文件名"},
                    "content": {"type": "string", "description": "要添加的内容"},
                    "heading_level": {"type": "integer", "description": "标题级别（0为普通段落，1-9为标题级别）", "default": 0}
                },
                "required": ["filename", "content"]
            }
        ),
        Tool(
            name="read_word_document",
            description="读取Word文档内容",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "文件名"}
                },
                "required": ["filename"]
            }
        ),
        
        # Excel工具
        Tool(
            name="create_excel_workbook",
            description="创建新的Excel工作簿",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "文件名（包含.xlsx扩展名）"},
                    "sheet_name": {"type": "string", "description": "工作表名称", "default": "Sheet1"}
                },
                "required": ["filename"]
            }
        ),
        Tool(
            name="write_excel_data",
            description="向Excel工作表写入数据",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "文件名"},
                    "sheet_name": {"type": "string", "description": "工作表名称"},
                    "data": {"type": "array", "description": "二维数组数据"},
                    "start_row": {"type": "integer", "description": "起始行", "default": 1},
                    "start_col": {"type": "integer", "description": "起始列", "default": 1}
                },
                "required": ["filename", "sheet_name", "data"]
            }
        ),
        Tool(
            name="read_excel_data",
            description="读取Excel工作表数据",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "文件名"},
                    "sheet_name": {"type": "string", "description": "工作表名称（可选）"}
                },
                "required": ["filename"]
            }
        ),
        
        # PowerPoint工具
        Tool(
            name="create_ppt_presentation",
            description="创建新的PPT演示文稿",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "文件名（包含.pptx扩展名）"},
                    "title": {"type": "string", "description": "演示文稿标题（可选）"}
                },
                "required": ["filename"]
            }
        ),
        Tool(
            name="add_ppt_slide",
            description="向PPT添加幻灯片",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "文件名"},
                    "title": {"type": "string", "description": "幻灯片标题"},
                    "content": {"type": "string", "description": "幻灯片内容（可选）"},
                    "layout_index": {"type": "integer", "description": "布局索引（默认1）", "default": 1}
                },
                "required": ["filename", "title"]
            }
        ),
        Tool(
            name="get_ppt_info",
            description="获取PPT演示文稿信息",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "文件名"}
                },
                "required": ["filename"]
            }
        ),
        
        # 通用工具
        Tool(
            name="list_files",
            description="列出工作区所有文件",
            inputSchema={
                "type": "object",
                "properties": {}
            }
        )
    ]

@server.call_tool()
async def handle_call_tool(name: str, arguments: dict) -> CallToolResult:
    """处理工具调用"""
    try:
        if name == "create_word_document":
            result = office_manager.create_word_document(
                arguments["filename"],
                arguments.get("title", ""),
                arguments.get("content", "")
            )
        elif name == "add_word_content":
            result = office_manager.add_word_content(
                arguments["filename"],
                arguments["content"],
                arguments.get("heading_level", 0)
            )
        elif name == "read_word_document":
            result = office_manager.read_word_document(arguments["filename"])
        elif name == "create_excel_workbook":
            result = office_manager.create_excel_workbook(
                arguments["filename"],
                arguments.get("sheet_name", "Sheet1")
            )
        elif name == "write_excel_data":
            result = office_manager.write_excel_data(
                arguments["filename"],
                arguments["sheet_name"],
                arguments["data"],
                arguments.get("start_row", 1),
                arguments.get("start_col", 1)
            )
        elif name == "read_excel_data":
            result = office_manager.read_excel_data(
                arguments["filename"],
                arguments.get("sheet_name")
            )
        elif name == "create_ppt_presentation":
            result = office_manager.create_ppt_presentation(
                arguments["filename"],
                arguments.get("title", "")
            )
        elif name == "add_ppt_slide":
            result = office_manager.add_ppt_slide(
                arguments["filename"],
                arguments["title"],
                arguments.get("content", ""),
                arguments.get("layout_index", 1)
            )
        elif name == "get_ppt_info":
            result = office_manager.get_ppt_info(arguments["filename"])
        elif name == "list_files":
            result = office_manager.list_files()
        else:
            result = f"未知工具: {name}"
        
        return CallToolResult(content=[TextContent(type="text", text=result)])
    
    except Exception as e:
        error_msg = f"执行工具 '{name}' 时出错: {str(e)}"
        return CallToolResult(content=[TextContent(type="text", text=error_msg)])

async def main():
    """启动MCP服务器"""
    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            InitializationOptions(
                server_name="office-manager",
                server_version="1.0.0",
                capabilities=server.get_capabilities(
                    notification_options=None,
                    experimental_capabilities=None
                )
            )
        )

if __name__ == "__main__":
    print("启动Office MCP服务器...")
    print("支持的功能:")
    print("- Word文档: 创建、编辑、读取")
    print("- Excel工作簿: 创建、数据读写")
    print("- PowerPoint演示文稿: 创建、添加幻灯片")
    print("- 文件管理: 列出工作区文件")
    
    asyncio.run(main())
